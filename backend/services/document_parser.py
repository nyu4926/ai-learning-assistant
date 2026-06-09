import os
import re
from typing import List, Tuple

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from docx import Document as DocxDocument

from config import CHUNK_SIZE, CHUNK_OVERLAP


# ── 格式映射 ────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".pdf":  "pdf",
    ".pptx": "pptx",
    ".ppt":  "ppt",
    ".docx": "docx",
    ".doc":  "doc",
    ".md":   "md",
    ".txt":  "txt",
}


def get_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}（支持: {', '.join(SUPPORTED_EXTENSIONS)}）")
    return SUPPORTED_EXTENSIONS[ext]


# ── 各格式解析器 ────────────────────────────────────────────────

def _parse_pdf(filepath: str) -> str:
    """PyMuPDF 提取 PDF 文本"""
    doc = fitz.open(filepath)
    parts: List[str] = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text")
        if text.strip():
            parts.append(f"[第 {page_num} 页]\n{text.strip()}")
    doc.close()
    return "\n\n".join(parts)


def _parse_pptx(filepath: str) -> str:
    """python-pptx 提取 PPTX 文本（含表格、幻灯片标头）"""
    prs = Presentation(filepath)
    parts: List[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slides_parts: List[str] = []
        for shape in slide.shapes:
            # 文本框
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slides_parts.append(t)
            # 表格
            if shape.has_table:
                table = shape.table
                rows_text: List[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(" | ".join(cells))
                slides_parts.append("\n".join(rows_text))
        if slides_parts:
            parts.append(f"## 幻灯片 {slide_num}\n" + "\n".join(slides_parts))
    return "\n\n".join(parts)


def _parse_ppt(filepath: str) -> str:
    """PPT（旧格式）尝试用 python-pptx 打开；部分新版支持"""
    # python-pptx 从 0.6.21 开始支持部分 .ppt
    return _parse_pptx(filepath)


def _parse_docx(filepath: str) -> str:
    """python-docx 提取 DOCX 文本（含表格）"""
    doc = DocxDocument(filepath)
    parts: List[str] = []
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
        if tag == "p":
            # 段落
            text = _extract_docx_paragraph_text(element)
            if text.strip():
                parts.append(text.strip())
        elif tag == "tbl":
            # 表格
            rows_text: List[str] = []
            for row in element.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tr"):
                cells = []
                for cell in row.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tc"):
                    cell_text = "".join(
                        node.text or ""
                        for node in cell.iter()
                        if node.tag.endswith("}t") and node.text
                    )
                    cells.append(cell_text.strip())
                rows_text.append(" | ".join(cells))
            parts.append("\n".join(rows_text))
    return "\n\n".join(parts)


def _extract_docx_paragraph_text(para_element) -> str:
    """从 docx paragraph XML 元素中提取纯文本"""
    texts = []
    for node in para_element.iter():
        if node.tag.endswith("}t") and node.text:
            texts.append(node.text)
    return "".join(texts)


def _parse_doc(filepath: str) -> str:
    """DOC（旧格式）：尝试用 python-docx 打开"""
    try:
        return _parse_docx(filepath)
    except Exception:
        raise ValueError(
            "无法解析 .doc 文件，请先转换为 .docx 格式"
        )


def _parse_markdown(filepath: str) -> str:
    """读取 Markdown 原始文本"""
    return _read_text_file(filepath)


def _parse_txt(filepath: str) -> str:
    """读取 TXT，自动处理编码"""
    return _read_text_file(filepath)


def _read_text_file(filepath: str) -> str:
    """尝试多种编码读取文本文件"""
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            with open(filepath, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法识别文件编码: {filepath}")


# ── 统一解析入口 ────────────────────────────────────────────────

PARSERS = {
    "pdf":  _parse_pdf,
    "pptx": _parse_pptx,
    "ppt":  _parse_ppt,
    "docx": _parse_docx,
    "doc":  _parse_doc,
    "md":   _parse_markdown,
    "txt":  _parse_txt,
}


def parse_document(filepath: str, file_type: str = None) -> str:
    """
    解析文档，返回纯文本。
    
    Args:
        filepath: 文件路径
        file_type: 格式类型（可从文件名推断）
    
    Returns:
        提取的纯文本字符串
    """
    if file_type is None:
        file_type = get_file_type(filepath)
    parser = PARSERS.get(file_type)
    if parser is None:
        raise ValueError(f"不支持的文档类型: {file_type}")
    return parser(filepath)


# ── 文本分块 ────────────────────────────────────────────────────

def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
    separator: str = "\n\n",
) -> List[Tuple[str, dict]]:
    """
    按语义段落切分文本。
    
    策略：
    1. 按 separator 拆分为段落
    2. 每个段落如果超过 chunk_size，递归按句子切分
    3. 合并短段落直到接近 chunk_size
    4. 相邻块之间保留 overlap 字符的重叠
    
    Returns:
        [(chunk_text, metadata_dict), ...]
    """
    if not text or not text.strip():
        return []

    paragraphs = [p.strip() for p in text.split(separator) if p.strip()]
    chunks: List[Tuple[str, dict]] = []
    current = ""
    
    for para in paragraphs:
        # 单个段落超过 chunk_size，递归按句子切分
        if len(para) > chunk_size:
            # 先保存当前累积的块
            if current:
                chunks.append((current, {"char_count": len(current)}))
                current = ""
            # 按句子切分长段落
            sub_chunks = _split_long_paragraph(para, chunk_size, overlap)
            chunks.extend(sub_chunks)
            continue

        # 合并短段落
        if current and len(current) + len(separator) + len(para) > chunk_size:
            chunks.append((current, {"char_count": len(current)}))
            # overlap：下一块的开头包含上一块的末尾
            if overlap > 0 and len(current) > overlap:
                current = current[-overlap:] + separator + para
            else:
                current = para
        else:
            if current:
                current += separator + para
            else:
                current = para

    if current:
        chunks.append((current, {"char_count": len(current)}))

    return chunks


def _split_by_sentence(text: str) -> List[str]:
    """按中英文句子标点切分"""
    # 保留标点在句尾
    parts = re.split(r"(?<=[。！？.!?\n])\s*", text)
    return [p.strip() for p in parts if p.strip()]


def _split_long_paragraph(
    text: str, chunk_size: int, overlap: int
) -> List[Tuple[str, dict]]:
    """将过长段落按句子切分"""
    sentences = _split_by_sentence(text)
    chunks: List[Tuple[str, dict]] = []
    current = ""

    for sent in sentences:
        if current and len(current) + len(sent) + 1 > chunk_size:
            chunks.append((current, {"char_count": len(current)}))
            if overlap > 0 and len(current) > overlap:
                current = current[-overlap:] + " " + sent
            else:
                current = sent
        else:
            if current:
                current += " " + sent
            else:
                current = sent

    if current:
        chunks.append((current, {"char_count": len(current)}))

    return chunks
