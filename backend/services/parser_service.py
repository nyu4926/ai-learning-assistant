"""文件解析服务 — 提取 PDF/PPT/Word/Markdown/TXT 文本内容"""

import os

import pdfplumber
from pptx import Presentation
from docx import Document


def extract_text(file_path: str) -> tuple[str, int]:
    """
    根据文件扩展名调用对应的解析器，返回 (纯文本, 页数)。
    解析失败时抛 ValueError。
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return _extract_ppt(file_path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(file_path)
    elif ext in (".md", ".markdown"):
        return _extract_text_file(file_path, "markdown")
    elif ext == ".txt":
        return _extract_text_file(file_path, "txt")
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _extract_pdf(file_path: str) -> tuple[str, int]:
    """用 pdfplumber 提取 PDF 文本"""
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    full_text = "\n\n".join(pages)
    return full_text, len(pages)


def _extract_ppt(file_path: str) -> tuple[str, int]:
    """用 python-pptx 提取 PPT 文本"""
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides_text.append("\n".join(texts))
    full_text = "\n\n".join(slides_text)
    return full_text, len(prs.slides)


def _extract_docx(file_path: str) -> tuple[str, int]:
    """用 python-docx 提取 Word 文本"""
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)
    # Word 没有"页数"概念，按段落数粗估
    page_count = max(1, len(paragraphs) // 20)
    return full_text, page_count


def _extract_text_file(file_path: str, file_type: str) -> tuple[str, int]:
    """Markdown / TXT 直接读取"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    line_count = text.count("\n") + 1
    return text, line_count
